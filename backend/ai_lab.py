"""AI Lab — free AI tools for founders.

Phase 1 tools:
  1. AI Pitch Deck Review — Gemini 2.5 Pro analyses the PDF and returns an
     8-dimension investor-readiness scorecard + top 3 priorities.
  2. Logo → Branded Template Generator — extracts brand colors from one (or
     two co-branded) logos and generates a 5-slide .pptx the user can open
     and edit. User can override any color via the color picker on the UI.

Both flows capture leads in the `ai_lab_leads` collection so SkiFi sees
the warm signal (active deck = active fundraise = ideal redesign client).
"""
from __future__ import annotations

import io
import json
import logging
import os
import re
import tempfile
import uuid
from datetime import datetime, timezone, timedelta
from typing import List, Optional

from colorthief import ColorThief
from fastapi import APIRouter, Cookie, Depends, File, Form, HTTPException, Request, Response, UploadFile
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from emergentintegrations.llm.chat import (
    FileContentWithMimeType,
    LlmChat,
    UserMessage,
)
logger = logging.getLogger(__name__)

ai_lab_router = APIRouter(prefix="/ai-lab", tags=["ai-lab"])

EMERGENT_LLM_KEY = os.environ.get("EMERGENT_LLM_KEY", "")
DECK_REVIEW_DAILY_LIMIT = int(os.environ.get("DECK_REVIEW_DAILY_LIMIT", "2"))
TEMPLATE_GEN_DAILY_LIMIT = int(os.environ.get("TEMPLATE_GEN_DAILY_LIMIT", "3"))

# Brand defaults - used when palette extraction fails or user skips picker
DEFAULT_PRIMARY = "#2A7AFE"
DEFAULT_DARK = "#0A0F1E"
DEFAULT_LIGHT = "#F5F4EE"
DEFAULT_TEXT = "#0A0A0A"

DECK_REVIEW_PROMPT = """You are SkiFi's AI deck reviewer, trained on what investors judge and on a studio that has shipped 2,700+ investor decks. Review the attached pitch deck. Respond with ONLY valid JSON, no markdown, no preamble, exactly this shape:
{"overall":<0-100 integer>,"verdict":"<=12 word investor-readiness verdict","dimensions":[{"name":"Narrative & Flow","score":<0-100>,"note":"<=16 words"},{"name":"Problem & Solution","score":0,"note":""},{"name":"Market Opportunity","score":0,"note":""},{"name":"Business Model","score":0,"note":""},{"name":"Traction & Metrics","score":0,"note":""},{"name":"Team & Credibility","score":0,"note":""},{"name":"Visual Design","score":0,"note":""},{"name":"The Ask","score":0,"note":""}],"priorities":["<=14 words","<=14 words","<=14 words"]}
Keep all 8 dimensions in that exact order. Be honest and specific."""


# ===================== Auth (mandatory Google sign-in) =====================
async def _require_signed_in_user(session_token: Optional[str]) -> dict:
    """Both AI Lab tools require Google sign-in. Returns the user dict or 401s."""
    from server import _get_user_from_token  # lazy to avoid circular imports
    user = await _get_user_from_token(session_token)
    if not user:
        raise HTTPException(
            status_code=401,
            detail="Please sign in with Google to use this free tool.",
        )
    return user


# ===================== Lead capture =====================
async def _record_lead(db, *, tool: str, name: str, email: str,
                        company: Optional[str], request: Request,
                        extra: Optional[dict] = None) -> str:
    """Stores a lead row and triggers an admin notification email."""
    lead_id = f"lead_{uuid.uuid4().hex[:12]}"
    doc = {
        "id": lead_id,
        "tool": tool,
        "name": name.strip(),
        "email": email.strip().lower(),
        "company": (company or "").strip() or None,
        "source_ip": request.headers.get("x-forwarded-for", "").split(",")[0].strip()
                     or request.client.host if request.client else "",
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    if extra:
        doc.update(extra)
    await db.ai_lab_leads.insert_one(doc)
    return lead_id


async def _enforce_daily_limit(db, *, tool: str, email: str, cap: int) -> None:
    cutoff = (datetime.now(timezone.utc) - timedelta(hours=24)).isoformat()
    used = await db.ai_lab_leads.count_documents({
        "tool": tool, "email": email.lower(), "created_at": {"$gte": cutoff}
    })
    if used >= cap:
        raise HTTPException(
            status_code=429,
            detail=f"Daily limit reached ({cap}/day). Reply to contact@skifidesigns.com and we'll review it personally.",
        )


# ===================== Tool 1: Deck Review =====================
def _parse_review_json(raw: str) -> dict:
    """LLMs sometimes wrap JSON in ```json blocks - strip & parse safely."""
    cleaned = raw.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    # Grab the first {...} block if there's stray text
    m = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if m:
        cleaned = m.group(0)
    return json.loads(cleaned)


@ai_lab_router.post("/deck-review")
async def deck_review(
    request: Request,
    company: str = Form(""),
    pdf: UploadFile = File(...),
    session_token: Optional[str] = Cookie(None),
):
    from server import db  # local import to avoid circular dep
    user = await _require_signed_in_user(session_token)
    name = user.get("name") or user.get("email", "").split("@")[0]
    email = user.get("email", "")
    if not EMERGENT_LLM_KEY:
        raise HTTPException(status_code=500, detail="LLM key not configured")
    if pdf.content_type not in ("application/pdf", "application/x-pdf"):
        raise HTTPException(status_code=400, detail="Please upload a PDF deck")

    # 20 MB safety cap
    body = await pdf.read()
    if len(body) > 20 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="Deck is larger than 20 MB")

    await _enforce_daily_limit(db, tool="deck_review", email=email,
                                cap=DECK_REVIEW_DAILY_LIMIT)

    # Persist to a temp file - the LLM lib reads paths, not bytes
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        f.write(body)
        tmp_path = f.name

    try:
        chat = (
            LlmChat(
                api_key=EMERGENT_LLM_KEY,
                session_id=f"deck-review-{uuid.uuid4().hex[:8]}",
                system_message="You return ONLY valid JSON. No markdown, no preamble.",
            )
            .with_model("gemini", "gemini-2.5-pro")
        )
        attachment = FileContentWithMimeType(file_path=tmp_path, mime_type="application/pdf")

        raw = await chat.send_message(UserMessage(
            text=DECK_REVIEW_PROMPT,
            file_contents=[attachment],
        ))
        raw = raw or ""
        try:
            result = _parse_review_json(raw)
        except Exception:
            logger.exception(f"Deck review JSON parse failed; raw={raw[:400]!r}")
            raise HTTPException(status_code=502,
                                detail="The AI returned an unreadable response - please try again")

        # Capture the lead (do this AFTER a successful review so failures don't burn the daily quota)
        await _record_lead(db, tool="deck_review", name=name, email=email,
                           company=company or None, request=request,
                           extra={"deck_filename": pdf.filename or "",
                                  "deck_size_bytes": len(body),
                                  "overall_score": result.get("overall"),
                                  "verdict": result.get("verdict")})
        return result
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


# ===================== Tool 2: Logo → Branded Template =====================
def _rgb_to_hex(rgb: tuple) -> str:
    return "#{:02X}{:02X}{:02X}".format(*[int(x) for x in rgb[:3]])


def _hex_to_rgb(hx: str) -> tuple:
    hx = (hx or "").lstrip("#")
    if len(hx) != 6:
        return (0, 0, 0)
    return tuple(int(hx[i:i + 2], 16) for i in (0, 2, 4))


def _luminance(rgb: tuple) -> float:
    """Rec. 709 relative luminance - higher = lighter."""
    r, g, b = [x / 255 for x in rgb[:3]]
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _extract_palette(img_bytes: bytes, color_count: int = 5) -> List[str]:
    """Returns a list of hex colors sorted with dominant first.
    Skips near-white/near-black so the palette feels brand-distinctive."""
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        f.write(img_bytes)
        tmp = f.name
    try:
        ct = ColorThief(tmp)
        dominant = ct.get_color(quality=1)
        # palette includes dominant; ask for a few more so we can filter
        palette = ct.get_palette(color_count=max(color_count + 2, 5), quality=1)
        seen = set()
        ordered = []
        for rgb in [dominant, *palette]:
            hx = _rgb_to_hex(rgb)
            if hx in seen:
                continue
            seen.add(hx)
            lum = _luminance(rgb)
            # Drop extreme white/black so we keep brand-distinctive colors
            if lum > 0.96 or lum < 0.04:
                continue
            ordered.append(hx)
            if len(ordered) >= color_count:
                break
        return ordered or ["#2A7AFE"]
    finally:
        try:
            os.unlink(tmp)
        except OSError:
            pass


@ai_lab_router.post("/template/extract-palette")
async def extract_palette(logo: UploadFile = File(...)):
    if not (logo.content_type or "").startswith("image/"):
        raise HTTPException(status_code=400, detail="Please upload an image (PNG/JPG/SVG)")
    body = await logo.read()
    if len(body) > 8 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="Logo larger than 8 MB")
    # Convert through Pillow to ensure ColorThief can read it (handles RGBA + WebP)
    try:
        img = Image.open(io.BytesIO(body)).convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        png_bytes = buf.getvalue()
    except Exception:
        raise HTTPException(status_code=400, detail="Could not read that image")
    palette = _extract_palette(png_bytes, color_count=5)
    primary = palette[0]
    accent = palette[1] if len(palette) > 1 else DEFAULT_PRIMARY
    return {"palette": palette, "primary": primary, "accent": accent}


# ===================== PPTX builder =====================
def _set_solid_fill(shape, hex_color: str) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*_hex_to_rgb(hex_color))


def _add_text(slide, *, left, top, width, height, text, color="#0A0A0A",
              size=18, bold=False, font="Outfit", align=None):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*_hex_to_rgb(color))
    return tb


def _build_pptx(*, logo1_bytes: Optional[bytes], logo2_bytes: Optional[bytes],
                primary: str, dark: str, light: str,
                project_name: str = "Your Brand",
                theme: str = "dark") -> bytes:
    """Builds a 5-slide branded template in either DARK or LIGHT theme.

    Slide list (same in both themes):
        1. Title slide (theme-colored bg + brand accent)
        2. Agenda
        3. Content (3-column)
        4. Data viz placeholder
        5. Closing / The Ask

    Theme behaviour:
        - dark  : title + closing on `dark`,  inner slides (2-4) on `light`
        - light : title + closing on `light`, inner slides (2-4) on `light`,
                  accents/headlines stay branded; primary stays brand-bright.
    """
    prs = Presentation()
    SLIDE_W = 13.333  # inches (16:9 widescreen)
    SLIDE_H = 7.5
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    blank_layout = prs.slide_layouts[6]

    is_dark = (theme or "dark").lower() == "dark"
    # Title/closing background tone toggles by theme. Inner content always on
    # the warm-light surface for readability — even in dark theme decks.
    hero_bg = dark if is_dark else light
    text_on_hero = "#FFFFFF" if _luminance(_hex_to_rgb(hero_bg)) < 0.5 else "#0A0A0A"
    inner_bg = light  # consistent inner surface
    text_on_inner = "#0A0A0A" if _luminance(_hex_to_rgb(inner_bg)) > 0.6 else "#FFFFFF"

    # ---------- Logo placement (top-right, never overflows the slide) ----------
    LOGO_W = 1.3
    LOGO_GAP = 0.18
    LOGO_TOP = 0.35
    LOGO_RIGHT_PAD = 0.5

    def add_logos(slide):
        """Top-right logos, right-aligned so a single OR a pair always
        sits inside the slide regardless of count."""
        logos = [b for b in (logo1_bytes, logo2_bytes) if b]
        if not logos:
            return
        n = len(logos)
        total_w = n * LOGO_W + (n - 1) * LOGO_GAP
        start_x = SLIDE_W - LOGO_RIGHT_PAD - total_w
        # Hard clamp so a wide logo never bleeds past the left margin
        if start_x < 0.4:
            start_x = 0.4
        cur_x = start_x
        for lb in logos:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(lb), Inches(cur_x), Inches(LOGO_TOP),
                    width=Inches(LOGO_W),
                )
            except Exception:
                pass
            cur_x += LOGO_W + LOGO_GAP

    def add_background(slide, hex_color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*_hex_to_rgb(hex_color))

    def add_accent_strip(slide, hex_color, x=0.6, y=0.6, w=1.5, h=0.07):
        rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        _set_solid_fill(rect, hex_color)
        rect.line.fill.background()

    # ---- 1. Title slide
    s1 = prs.slides.add_slide(blank_layout)
    add_background(s1, hero_bg)
    add_logos(s1)
    add_accent_strip(s1, primary, x=0.7, y=2.6, w=1.4, h=0.08)
    _add_text(s1, left=0.7, top=2.8, width=11, height=1.4,
              text=project_name, color=text_on_hero, size=60, bold=True, font="Nohemi")
    _add_text(s1, left=0.7, top=4.3, width=11, height=0.6,
              text="Tagline goes here", color=text_on_hero, size=22, font="Outfit")
    _add_text(s1, left=0.7, top=6.6, width=8, height=0.4,
              text="Presented by — Date", color=text_on_hero, size=12, font="Outfit")

    # ---- 2. Agenda
    s2 = prs.slides.add_slide(blank_layout)
    add_background(s2, inner_bg)
    add_logos(s2)
    _add_text(s2, left=0.7, top=0.6, width=8, height=0.5,
              text="AGENDA", color=primary, size=14, bold=True, font="Outfit")
    _add_text(s2, left=0.7, top=1.0, width=10, height=1.2,
              text="What we'll cover", color=text_on_inner, size=44, bold=True, font="Nohemi")
    agenda = ["01  The opportunity", "02  Our solution", "03  Traction & roadmap",
              "04  Team", "05  The ask"]
    for i, item in enumerate(agenda):
        _add_text(s2, left=0.7, top=2.8 + i * 0.7, width=10, height=0.55,
                  text=item, color=text_on_inner, size=22, font="Outfit")

    # ---- 3. Three-column content
    s3 = prs.slides.add_slide(blank_layout)
    add_background(s3, inner_bg)
    add_logos(s3)
    _add_text(s3, left=0.7, top=0.6, width=8, height=0.5,
              text="OUR SOLUTION", color=primary, size=14, bold=True, font="Outfit")
    _add_text(s3, left=0.7, top=1.0, width=10, height=1.2,
              text="Three pillars", color=text_on_inner, size=44, bold=True, font="Nohemi")
    cols = [
        ("01", "Pillar one", "A short, punchy description of the first value pillar - what it does and why it matters."),
        ("02", "Pillar two", "A short, punchy description of the second value pillar - the unique edge you bring."),
        ("03", "Pillar three", "A short, punchy description of the third value pillar - what wraps it together."),
    ]
    for i, (num, title, desc) in enumerate(cols):
        x = 0.7 + i * 4.2
        rect = s3.shapes.add_shape(1, Inches(x), Inches(2.9), Inches(0.6), Inches(0.6))
        _set_solid_fill(rect, primary)
        rect.line.fill.background()
        _add_text(s3, left=x, top=2.95, width=0.6, height=0.5, text=num,
                  color="#FFFFFF", size=20, bold=True, font="Nohemi", align="center")
        _add_text(s3, left=x, top=3.7, width=3.8, height=0.6, text=title,
                  color=text_on_inner, size=24, bold=True, font="Nohemi")
        _add_text(s3, left=x, top=4.4, width=3.8, height=2.5, text=desc,
                  color="#555555", size=14, font="Outfit")

    # ---- 4. Data viz placeholder
    s4 = prs.slides.add_slide(blank_layout)
    add_background(s4, inner_bg)
    add_logos(s4)
    _add_text(s4, left=0.7, top=0.6, width=8, height=0.5,
              text="TRACTION", color=primary, size=14, bold=True, font="Outfit")
    _add_text(s4, left=0.7, top=1.0, width=10, height=1.2,
              text="The numbers", color=text_on_inner, size=44, bold=True, font="Nohemi")
    metrics = [("12K+", "Active users"), ("$420K", "ARR"), ("4.8x", "YoY growth")]
    for i, (val, lbl) in enumerate(metrics):
        x = 0.7 + i * 4.2
        rect = s4.shapes.add_shape(5, Inches(x), Inches(3.2), Inches(3.8), Inches(2.8))
        _set_solid_fill(rect, "#FFFFFF")
        rect.line.color.rgb = RGBColor(*_hex_to_rgb("#ECEAE2"))
        _add_text(s4, left=x + 0.3, top=3.5, width=3.5, height=1.2, text=val,
                  color=primary, size=48, bold=True, font="Nohemi")
        _add_text(s4, left=x + 0.3, top=4.9, width=3.5, height=0.6, text=lbl,
                  color="#555555", size=16, font="Outfit")

    # ---- 5. Closing / Ask
    s5 = prs.slides.add_slide(blank_layout)
    add_background(s5, hero_bg)
    add_logos(s5)
    add_accent_strip(s5, primary, x=0.7, y=2.6, w=1.4, h=0.08)
    _add_text(s5, left=0.7, top=2.8, width=12, height=1.6,
              text="Let's build the future, together.", color=text_on_hero, size=52,
              bold=True, font="Nohemi")
    _add_text(s5, left=0.7, top=5.0, width=11, height=0.7,
              text="hello@yourbrand.com  -  yourbrand.com",
              color=text_on_hero, size=18, font="Outfit")
    _add_text(s5, left=0.7, top=6.6, width=10, height=0.4,
              text="Designed by SkiFi Designs - skifidesigns.com",
              color=text_on_hero, size=11, font="Outfit")

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


@ai_lab_router.post("/template/generate")
async def generate_template(
    request: Request,
    company: str = Form(""),
    project_name: str = Form("Your Brand"),
    primary: str = Form(DEFAULT_PRIMARY),
    dark: str = Form(DEFAULT_DARK),
    light: str = Form(DEFAULT_LIGHT),
    theme: str = Form("dark"),
    logo1: UploadFile = File(...),
    logo2: Optional[UploadFile] = File(None),
    session_token: Optional[str] = Cookie(None),
):
    from server import db
    user = await _require_signed_in_user(session_token)
    name = user.get("name") or user.get("email", "").split("@")[0]
    email = user.get("email", "")
    theme_norm = theme.lower().strip() if theme else "dark"
    if theme_norm not in ("dark", "light"):
        theme_norm = "dark"
    await _enforce_daily_limit(db, tool="template_generator", email=email,
                                cap=TEMPLATE_GEN_DAILY_LIMIT)

    logo1_bytes = await logo1.read()
    logo2_bytes = None
    if logo2 and logo2.filename:
        logo2_bytes = await logo2.read()

    # Normalise to PNG so python-pptx renders them cleanly (handles SVG/WebP poorly)
    def _to_png(raw: bytes) -> bytes:
        try:
            img = Image.open(io.BytesIO(raw))
            img = img.convert("RGBA") if img.mode in ("P", "LA", "RGBA") else img.convert("RGB")
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()
        except Exception:
            return raw

    pptx_bytes = _build_pptx(
        logo1_bytes=_to_png(logo1_bytes),
        logo2_bytes=_to_png(logo2_bytes) if logo2_bytes else None,
        primary=primary, dark=dark, light=light,
        project_name=project_name,
        theme=theme_norm,
    )

    await _record_lead(db, tool="template_generator", name=name, email=email,
                       company=company or None, request=request,
                       extra={"project_name": project_name, "primary": primary,
                              "dark": dark, "light": light, "theme": theme_norm,
                              "co_branded": logo2_bytes is not None})

    safe_name = re.sub(r"[^A-Za-z0-9_-]+", "_", project_name)[:40] or "SkiFi-Template"
    filename = f"SkiFi-{safe_name}-Template.pptx"
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ===================== Admin =====================
def register_admin_routes(api_router, require_admin):
    """Adds admin routes to the main api_router so they sit under /api/admin/*."""

    @api_router.get("/admin/ai-lab/leads")
    async def admin_list_leads(
        tool: Optional[str] = None,
        _: str = Depends(require_admin),
    ):
        from server import db
        q = {}
        if tool:
            q["tool"] = tool
        items = await db.ai_lab_leads.find(q, {"_id": 0}).sort("created_at", -1).to_list(500)
        return {"items": items}

    @api_router.get("/admin/ai-lab/leads.csv")
    async def admin_export_leads(_: str = Depends(require_admin)):
        from server import db
        import csv
        items = await db.ai_lab_leads.find({}, {"_id": 0}).sort("created_at", -1).to_list(5000)
        buf = io.StringIO()
        writer = csv.DictWriter(buf, fieldnames=[
            "created_at", "tool", "name", "email", "company", "overall_score",
            "verdict", "deck_filename", "project_name", "primary", "dark", "light",
            "theme", "co_branded", "source_ip",
        ], extrasaction="ignore")
        writer.writeheader()
        for it in items:
            writer.writerow(it)
        return Response(
            content=buf.getvalue(),
            media_type="text/csv",
            headers={"Content-Disposition": 'attachment; filename="ai-lab-leads.csv"'},
        )
